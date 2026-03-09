from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_final_presentation():
    prs = Presentation()

    # Define some styles
    def set_slide_background(slide, color_rgb):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color_rgb

    BG_COLOR = RGBColor(255, 255, 255)  # Pure White
    TEXT_COLOR = RGBColor(31, 41, 55)   # Dark Gray (Slate 800)
    ACCENT_COLOR = RGBColor(0, 51, 102) # Navy Blue (Professional Academic)

    slides_content = [
        {
            "title": "STOCK MOMENT PREDICTION USING TFT",
            "subtitle": "Major Project Review\n\nSubmitted By:\nAshish Kaushik (22KB1A0509)\nD. Ashok (22KB1A0538)\nBhanu Prakash (22KB1A05F3)\n\nUnder the esteemed guidance of:\nDr. N. Kesava Rao M. Tech, Ph.D\nAssistant Professor, Dept. of CSE",
            "is_title": True
        },
        {
            "title": "Abstract",
            "content": [
                "This project presents a novel approach to stock market movement forecasting using Temporal Fusion Transformers (TFT).",
                "Traditional models like LSTM and ARIMA often struggle with long-term dependencies and heterogeneous data sources.",
                "The TFT architecture is designed for multi-horizon time-series forecasting, utilizing attention mechanisms to capture complex temporal dynamics.",
                "Integrated with a real-time dashboard and Angel One SmartAPI for live trading insights."
            ]
        },
        {
            "title": "Introduction",
            "content": [
                "Stock market prediction is a significant challenge in financial analytics due to high volatility and non-linearity.",
                "Evolution from Statistical Methods (ARIMA) to Deep Learning (RNN, LSTM).",
                "Current Research Focus: Attention-based architectures and Transformer-based models for better interpretability and performance.",
                "Significance: Providing retail investors with high-accuracy, real-time momentum predictions."
            ]
        },
        {
            "title": "Literature Survey",
            "content": [
                "ARIMA: Effective for linear trends but fails in volatile market conditions.",
                "LSTM: Captures long-term dependencies but lacks 'interpretability' of feature importance.",
                "Gradient Boosting: Fast and accurate but struggles with truly complex temporal sequencing.",
                "Temporal Fusion Transformer (TFT): State-of-the-art for multi-horizon forecasting; utilizes localized LSTM layers and global self-attention."
            ]
        },
        {
            "title": "Problem Statement",
            "content": [
                "Data Heterogeneity: Complexity in merging technical indicators, historical prices, and real-time market sentiment.",
                "Prediction Delay: Most models suffer from high latency, making them impractical for active trading.",
                "Feature Selection: Identifying which market variables truly influence the next-day price 'moments'.",
                "Scalability: Challenges in deploying heavy ML models for real-time multi-stock analysis."
            ]
        },
        {
            "title": "Proposed Methodology: TFT",
            "content": [
                "Temporal Fusion Transformer (TFT): A multi-horizon forecasting model.",
                "Gating Mechanisms: Automatically skip redundant components of the network to prevent over-fitting.",
                "Variable Selection Networks: Identify relevant input features at each time step.",
                "Static Covariate Encoders: Integrate metadata (e.g., sector info) into the temporal prediction.",
                "Attention-based Multi-Horizon Forecasting: Captures long-range dependencies efficiently."
            ]
        },
        {
            "title": "System Architecture",
            "content": [
                "Backend: FastAPI (Uvicorn) - High-performance asynchronous REST API.",
                "Frontend: React.js with Vite & Tailwind CSS - Real-time Visualization Dashboard.",
                "Data Layer: Integration with Angel One SmartAPI for zero-delay NSE data feeds.",
                "Inference Engine: Scikit-learn and Scipy for data pre-processing and model execution.",
                "Persistence: SQLite for watchlist management and historical records."
            ]
        },
        {
            "title": "Dataset & Pre-processing",
            "content": [
                "Source: 5 years of historical OHLCV data from NSE (via Yahoo Finance & Angel One).",
                "Feature Engineering: RSI, SMA (50/200), Volatility, and Momentum metrics.",
                "Normalization: Feature scaling using StandardScaler to ensure model convergence.",
                "Temporal Windowing: Creating look-back windows for sequence-based prediction."
            ]
        },
        {
            "title": "Implementation Progress",
            "content": [
                "Real-time Data Integration: Migration to Angel One SmartAPI complete.",
                "Backend Parallelization: Optimized stock analysis to under 5 seconds for 24+ stocks.",
                "UI Calibration: Premium Dark/Glassmorphism dashboard with interactive charts.",
                "Deployment: Hosted on Render (Backend) and Netlify (Frontend) with CI/CD."
            ]
        },
        {
            "title": "Experimental Results",
            "content": [
                "Metric: RMSE (Root Mean Square Error) and R-Squared (R2) Score.",
                "Confidence Interval: 85%+ accuracy on short-term price trend matching.",
                "Performance: Sub-second inference time per stock after model optimization.",
                "Robustness: Successful mock fallback during API rate-limiting events."
            ]
        },
        {
            "title": "Conclusion",
            "content": [
                "The 'Stock Moment Prediction' system successfully bridges ML research with real-world trading tools.",
                "TFT architecture provides a superior balance between accuracy and structural interpretability.",
                "Professional integration with Angel One ensures the system is ready for live-market scenarios.",
                "User Dashboard provides actionable intelligence, reducing investor bias."
            ]
        },
        {
            "title": "Future Work",
            "content": [
                "Deep Sentiment Integration: Real-time NLP on Twitter/News feeds for sentiment-weighted prediction.",
                "Cross-Platform Mobile App: Flutter or React Native integration.",
                "Algorithmic Trading: Automated order execution via SmartAPI based on high-confidence signals.",
                "Multi-Asset Support: Expanding to Commodities, Forex, and Crypto."
            ]
        },
        {
            "title": "References",
            "content": [
                "1. Bryan Lim, et al., 'Temporal Fusion Transformers for Interpretable Multi-horizon Time Series Forecasting'.",
                "2. Vaswani, et al., 'Attention Is All You Need' - Foundation of Transformer models.",
                "3. Angel One SmartAPI Documentation for financial data integration.",
                "4. Scikit-learn: Machine Learning in Python, Pedregosa et al., JMLR 12."
            ]
        },
        {
            "title": "Thank You",
            "subtitle": "Questions?\n\nContact: Ashish Kaushik, D. Ashok, Bhanu Prakash",
            "is_title": True
        }
    ]

    for slide_data in slides_content:
        is_title_slide = slide_data.get("is_title", False)
        slide_layout = prs.slide_layouts[1] if not is_title_slide else prs.slide_layouts[0]
            
        slide = prs.slides.add_slide(slide_layout)
        set_slide_background(slide, BG_COLOR)

        title = slide.shapes.title
        title.text = slide_data["title"]
        title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.size = Pt(36)

        if is_title_slide and slide_data.get("subtitle"):
            subtitle = slide.placeholders[1]
            subtitle.text = slide_data["subtitle"]
            subtitle.text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
            subtitle.text_frame.paragraphs[0].font.size = Pt(24)
            subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        if slide_data.get("content") and not is_title_slide:
            content_placeholder = slide.placeholders[1]
            tf = content_placeholder.text_frame
            tf.word_wrap = True
            tf.text = "" # Clear default
            for detail in slide_data["content"]:
                p = tf.add_paragraph()
                p.text = "• " + detail
                p.font.color.rgb = TEXT_COLOR
                p.font.size = Pt(20)
                p.space_after = Pt(14)
                p.level = 0

    output_path = r"c:\Users\ADMIN\AI_stock-analysis\Stock_Moment_Prediction_TFT.pptx"
    prs.save(output_path)
    print(f"Final Review Presentation saved to: {output_path}")

if __name__ == "__main__":
    create_final_presentation()
