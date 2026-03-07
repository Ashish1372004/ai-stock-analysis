from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_presentation():
    prs = Presentation()

    # Define some styles
    def set_slide_background(slide, color_rgb):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color_rgb

    BG_COLOR = RGBColor(255, 255, 255)  # Pure White
    TEXT_COLOR = RGBColor(31, 41, 55)   # Dark Gray (Slate 800)
    ACCENT_COLOR = RGBColor(5, 150, 105) # Emerald 600 (Professional Green)

    slides_content = [
        {
            "title": "Project: AI-Stock Market Analysis Dashboard",
            "subtitle": "Advanced Financial Intelligence & Machine Learning Integration\nA Comprehensive Technical Overview",
            "content": []
        },
        {
            "title": "1. Project Vision & Problem Statement",
            "content": [
                "Objective: To bridge the gap between complex quantitative analysis and individual retail investors.",
                "Target Market: Focused on Indian equities (NSE/BSE) using real-time Yahoo Finance integration.",
                "Problem Solved: Eliminates 'Analysis Paralysis' by condensing complex technical indicators into clear Buy/Hold/Sell signals.",
                "Value Prop: High-fidelity financial charts combined with modern UX/UI standards for a professional trading experience."
            ]
        },
        {
            "title": "2. Advanced Technical Analysis Engine",
            "content": [
                "Automated Moving Average Crossovers: Monitoring 50-day SMA vs 200-day SMA for robust trend identification.",
                "Momentum Tracking: Real-time RSI (Relative Strength Index) calculation with dynamic oversold/overbought thresholds.",
                "Volumetric Analysis: Assessing volume spikes to confirm price breakout validity.",
                "Persistence Layer: SQLite implementation for long-term user watchlist and alert configuration storage.",
                "Dynamic News Aggregation: Tailored 'Market Intel' feed filtered by specific stock relevance using customized metadata extraction."
            ]
        },
        {
            "title": "3. Machine Learning Architecture",
            "content": [
                "Algorithm: Implementation of the Random Forest Regressor for multi-variable price prediction.",
                "Dataset: Trained on 5 years of historical OHLCV data fetched via yfinance asynchronously.",
                "Feature Engineering: Custom extraction of Volatility, 5/20-day Momentum, and Moving Average convergence metrics.",
                "Inference Pipeline: Seamless hand-off from data pre-processing to model inference in the FastAPI backend.",
                "Confidence Logic: Quantifying the reliability of AI predictions based on model variance and data completeness."
            ]
        },
        {
            "title": "4. Premium Front-End Design & UX",
            "content": [
                "Design Philosophy: Glassmorphism and Ultra-Dark aesthetics utilizing Tailwind CSS and Vanilla CSS variables.",
                "Interactive Elements: Custom Canvas-based Particle Mesh cursor trail (Spring physics & sine wave movements).",
                "Advanced Search: Intelligent <StockSearch /> component with trie-based suggestion logic for NSE symbols.",
                "Grid Pagination: Efficient 'Load More' logic managing dashboard performance and initial page load speed.",
                "Animation Orchestration: Framer Motion utilized for staggered entrance animations and smooth view transitions."
            ]
        },
        {
            "title": "5. System Architecture & Infrastructure",
            "content": [
                "Frontend: React.js with Vite for optimized builds and lightning-fast Hot Module Replacement.",
                "Backend: FastAPI (Uvicorn) providing a high-performance, asynchronous REST API layer.",
                "CI/CD: Automated deployment pipelines configured at the root via Netlify and Render.",
                "Routing & Proxy: Configured Netlify redirects for production-level API proxying to bypass CORS issues.",
                "Robustness: Global error handling with mock fallback data to ensure high availability during API rate limits."
            ]
        },
        {
            "title": "6. Implementation Milestones",
            "content": [
                "Phase 1: Core Technical Analysis engine and React-based dashboard visualization.",
                "Phase 2: Integration of Scikit-Learn based Machine Learning models for predictive analytics.",
                "Phase 3: UI Enhancement cycle: Custom cursor effects, search auto-suggestions, and dashboard pagination.",
                "Phase 4: Production-ready cleanup: TypeScript strict build verification and Netlify deployment optimization."
            ]
        },
        {
            "title": "7. Summary & Future Outlook",
            "content": [
                "Current State: Fully functional MVP with integrated ML inference and premium analytics UI.",
                "Future Roadmap: Real-time WebSocket streaming for live price ticks (Nifty 50).",
                "Future Feature: Cross-platform mobile integration using React Native.",
                "Future Feature: Sentiment analysis on social media/news feeds using NLP (Natural Language Processing)."
            ]
        }
    ]

    for slide_data in slides_content:
        slide_layout = prs.slide_layouts[1] # Title and Content
        if slide_data["title"].startswith("Project:"):
            slide_layout = prs.slide_layouts[0] # Title Slide
            
        slide = prs.slides.add_slide(slide_layout)
        set_slide_background(slide, BG_COLOR)

        title = slide.shapes.title
        title.text = slide_data["title"]
        title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
        title.text_frame.paragraphs[0].font.bold = True

        if slide_data.get("subtitle") and slide_layout == prs.slide_layouts[0]:
            subtitle = slide.placeholders[1]
            subtitle.text = slide_data["subtitle"]
            subtitle.text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
        
        if slide_data.get("content") and slide_layout != prs.slide_layouts[0]:
            content_placeholder = slide.placeholders[1]
            tf = content_placeholder.text_frame
            tf.word_wrap = True
            tf.text = "" # Clear default
            for detail in slide_data["content"]:
                p = tf.add_paragraph()
                p.text = "• " + detail
                p.font.color.rgb = TEXT_COLOR
                p.font.size = Pt(18)
                p.space_after = Pt(12)
                p.level = 0

    output_path = r"c:\Users\ADMIN\AI_stock-analysis\AI_Stock_Project_Detailed.pptx"
    prs.save(output_path)
    print(f"Detailed Presentation saved to: {output_path}")

if __name__ == "__main__":
    create_presentation()
